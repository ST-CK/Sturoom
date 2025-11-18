from fastapi import APIRouter, Request, Header, HTTPException
from fastapi.responses import JSONResponse
import os, json, httpx, requests
from io import BytesIO
from datetime import datetime, timedelta, timezone
from typing import Tuple
from dotenv import load_dotenv
from openai import OpenAI
from PyPDF2 import PdfReader
from pptx import Presentation
from supabase import create_client, Client
from pathlib import Path

# ---------------- 초기 설정 ----------------
load_dotenv(dotenv_path=Path(__file__).resolve().parent.parent / ".env")

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")

if not (OPENAI_API_KEY and SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY):
    raise RuntimeError("환경변수 누락: OPENAI_API_KEY / SUPABASE_URL / SUPABASE_SERVICE_ROLE_KEY")

client = OpenAI(api_key=OPENAI_API_KEY)
supabase: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)
router = APIRouter()
MAX_TOTAL_CHARS = 18000
KST = timezone(timedelta(hours=9))

# ---------------- 유틸 ----------------
def _safe_cut(text: str, limit: int) -> str:
    return (text or "").strip()[:limit]

async def _download_file(url: str) -> Tuple[str, bytes]:
    async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as http:
        r = await http.get(url)
    if r.status_code != 200:
        raise RuntimeError(f"다운로드 실패({r.status_code})")
    return url.split("/")[-1].lower(), r.content

def _extract_text_from_pdf(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    return "\n".join([page.extract_text() or "" for page in reader.pages])

def _extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    slides = []
    for slide in prs.slides:
        texts = [s.text for s in slide.shapes if hasattr(s, "text") and s.text]
        slides.append("\n".join(texts))
    return "\n".join(slides)

def _build_prompt(all_text: str, mode: str) -> str:
    mode_map = {
        "ox": "OX 형식 문제 (정답은 O 또는 X)",
        "short": "서술형 문제 (짧은 한 문장으로 답변)",
        "multiple": "객관식 4지선다 문제",
        "mixed": "객관식 4지선다, OX, 서술형을 섞은 혼합형 문제",
    }
    return f"""
다음은 강의 자료의 통합 텍스트입니다.
이 내용을 바탕으로 학습 이해도를 평가할 수 있는 {mode_map.get(mode, "혼합형 문제")} 3문항을 만들어 주세요.

요구사항:
1) JSON 배열만 출력
2) 각 문항: question, choices(보기 리스트), answer(정답), explanation(해설)
3) 보기 앞에 'A.', 'B.' 같은 접두사 금지
4) JSON 외의 텍스트 출력 금지

-----
{_safe_cut(all_text, MAX_TOTAL_CHARS)}
-----
"""

# ---------------- Supabase 토큰 검증 ----------------
def verify_supabase_token(token: str):
    url = f"{SUPABASE_URL}/auth/v1/user"
    headers = {
        "Authorization": f"Bearer {token}",
        "apikey": SUPABASE_SERVICE_ROLE_KEY,  # Render에서도 동작하게
    }
    try:
        res = requests.get(url, headers=headers, timeout=10)
        if res.status_code == 200:
            print("✅ Supabase 인증 성공")
            return res.json()
        print(f"❌ Supabase 인증 실패: {res.status_code} {res.text}")
        raise HTTPException(status_code=401, detail="Supabase 인증 실패")
    except requests.exceptions.RequestException as e:
        print("🚨 Supabase API 연결 실패:", e)
        raise HTTPException(status_code=500, detail="Supabase 연결 실패")

# ---------------- 세션 & 실행(run) 생성 (항상 새로운 세션) ----------------
@router.post("/session/start")
async def start_quiz_session(req: Request, authorization: str = Header(None)):
    data = await req.json()
    room_id = data.get("room_id")
    week_id = data.get("post_id")
    mode = data.get("mode", "mixed")

    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="인증 토큰이 없습니다.")

    token = authorization.split(" ")[1]

    # 🔐 유저 인증
    try:
        user_data = verify_supabase_token(token)
        user_id = user_data["id"]
        print(f"🧩 인증된 사용자: {user_id}")
    except:
        raise

    try:
        # ❗ 항상 새로운 세션 생성
        s_res = supabase.table("quiz_sessions").insert({
            "user_id": user_id,
            "lecture_id": room_id,
            "week_id": week_id,
            "mode": mode,
            "quiz_count": 0,
            "created_at": datetime.now(KST).isoformat()
        }).execute()

        if not s_res.data:
            raise RuntimeError("세션 생성 실패")

        session_id = s_res.data[0]["id"]

        # 새 run 생성
        r_res = supabase.table("quiz_runs").insert({
            "session_id": session_id,
            "user_id": user_id,
            "lecture_id": room_id,
            "week_id": week_id,
            "mode": mode,
            "started_at": datetime.now(KST).isoformat()
        }).execute()

        if not r_res.data:
            raise RuntimeError("런 생성 실패")

        run_id = r_res.data[0]["id"]

        print(f"🆕 새 세션/런 생성 완료: session={session_id}, run={run_id}")
        return JSONResponse({"session_id": session_id, "run_id": run_id})

    except Exception as e:
        print("❌ 세션 생성 오류:", repr(e))
        return JSONResponse(status_code=500, content={"error": str(e)})

# ---------------- 실행(run)만 생성 (기존 세션 재도전용) ----------------
@router.post("/run/start")
async def start_quiz_run(req: Request, authorization: str = Header(None)):
    """
    기존 session_id를 받아서 그 세션에 속한 새 run만 생성.
    - session_id는 그대로
    - quiz_runs에만 새로운 row 추가
    """
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="인증 토큰이 없습니다.")

    token = authorization.split(" ")[1]

    # 🔐 유저 인증
    try:
        user_data = verify_supabase_token(token)
        user_id = user_data["id"]
        print(f"🧩 인증된 사용자(재도전 run): {user_id}")
    except:
        raise

    try:
        data = await req.json()
        session_id = data.get("session_id")

        if not session_id:
            return JSONResponse(
                status_code=400,
                content={"error": "session_id가 필요합니다."},
            )

        # 기존 세션 정보 조회 (lecture_id, week_id, mode 재사용)
        s_res = (
            supabase.table("quiz_sessions")
            .select("lecture_id, week_id, mode")
            .eq("id", session_id)
            .limit(1)
            .execute()
        )

        if not s_res.data:
            return JSONResponse(
                status_code=404,
                content={"error": "해당 session_id의 세션이 존재하지 않습니다."},
            )

        session = s_res.data[0]
        lecture_id = session.get("lecture_id")
        week_id = session.get("week_id")
        mode = session.get("mode") or "mixed"

        # 새 run 생성
        r_res = (
            supabase.table("quiz_runs")
            .insert({
                "session_id": session_id,
                "user_id": user_id,
                "lecture_id": lecture_id,
                "week_id": week_id,
                "mode": mode,
                "quiz_count": 0,
                "started_at": datetime.now(KST).isoformat(),
            })
            .execute()
        )

        if not r_res.data:
            raise RuntimeError("런 생성 실패")

        run_id = r_res.data[0]["id"]
        print(f"🔁 기존 세션에 새 런 생성: session={session_id}, run={run_id}")

        return JSONResponse({"session_id": session_id, "run_id": run_id})

    except Exception as e:
        print("❌ 런 생성 오류:", repr(e))
        return JSONResponse(status_code=500, content={"error": str(e)})

# ---------------- 퀴즈 생성 ----------------
@router.post("/from-url")
async def generate_quiz_from_url(req: Request):
    data = await req.json()
    file_urls = data.get("file_urls") or []
    mode = (data.get("mode") or "mixed").strip().lower()
    user_id = data.get("user_id")
    room_id = data.get("room_id")
    week_id = data.get("week_id")
    session_id = data.get("session_id")
    run_id = data.get("run_id")

    if not file_urls:
        return JSONResponse(status_code=400, content={"error": "file_urls가 없습니다."})

    # 파일 내용 합치기
    aggregated = []
    for file in file_urls:
        try:
            url = file.get("url") if isinstance(file, dict) else file
            fname, blob = await _download_file(url)
            text = _extract_text_from_pdf(blob) if fname.endswith(".pdf") else _extract_text_from_pptx(blob)
            aggregated.append(f"\n### {fname}\n{text}")
        except Exception as e:
            print(f"⚠ 파일 처리 실패: {url} ({e})")

    all_text = "\n".join(aggregated)
    prompt = _build_prompt(all_text, mode)

    # AI 호출
    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "너는 교육용 퀴즈를 JSON으로만 반환하는 AI 교사야."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.2,
        )
        quiz_text = resp.choices[0].message.content
        json_str = quiz_text[quiz_text.find("["):quiz_text.rfind("]") + 1]
        quiz_data = json.loads(json_str)
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": f"OpenAI 처리 실패: {str(e)}"})

    # Supabase 저장
    try:
        questions = []
        for q in quiz_data:
            answer = (q.get("answer") or q.get("correct_answer") or "").strip()
            explanation = (q.get("explanation") or "").strip()
            choices = q.get("choices") or q.get("options") or []

            cleaned = []
            for c in choices:
                if isinstance(c, str):
                    cleaned.append(
                        c.replace("A. ", "")
                         .replace("B. ", "")
                         .replace("C. ", "")
                         .replace("D. ", "")
                         .strip()
                    )
                else:
                    cleaned.append(c)

            questions.append({
                "session_id": session_id,
                "question": q.get("question", "").strip(),
                "choices": cleaned,
                "answer": answer,
                "explanation": explanation,
            })

        inserted = supabase.table("quiz_questions").insert(questions).execute()

        # 세션/런 quiz_count 업데이트
        supabase.table("quiz_runs").update({"quiz_count": len(inserted.data)}).eq("id", run_id).execute()
        supabase.table("quiz_sessions").update({"quiz_count": len(inserted.data)}).eq("id", session_id).execute()

        # 첫 문제 메시지 저장 (전체 퀴즈 목록 payload로)
        supabase.table("quiz_messages").insert({
            "session_id": session_id,
            "run_id": run_id,
            "user_id": user_id,
            "role": "ai",
            "kind": "quiz",
            "payload": json.dumps({"quiz": inserted.data}),
        }).execute()

        return JSONResponse({
            "message": "퀴즈 생성 완료",
            "session_id": session_id,
            "run_id": run_id,
            "quiz_count": len(inserted.data),
            "quiz": inserted.data
        })

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

# ---------------- 정답 채점 ----------------
@router.post("/attempt")
async def attempt(req: Request):
    payload = await req.json()
    session_id = payload.get("session_id")
    run_id = payload.get("run_id")
    question_id = payload.get("question_id")
    user_email = payload.get("user_email")
    user_answer = (payload.get("user_answer") or "").strip()

    try:
        # 유저 찾기
        user_res = supabase.table("profiles").select("id").eq("email", user_email).limit(1).execute()
        if not user_res.data:
            return JSONResponse(status_code=404, content={"error": "유저 없음"})
        user_id = user_res.data[0]["id"]

        # 정답 + 해설 로드
        q_res = (
            supabase.table("quiz_questions")
            .select("answer, explanation")
            .eq("id", question_id)
            .limit(1)
            .execute()
        )
        if not q_res.data:
            return JSONResponse(status_code=404, content={"error": "문항 없음"})

        correct_answer = (q_res.data[0]["answer"] or "").strip()
        explanation = (q_res.data[0].get("explanation") or "").strip()

        # 채점
        is_correct = user_answer.lower() == correct_answer.lower()

        # 답안 저장
        supabase.table("quiz_answers").insert({
            "user_id": user_id,
            "question_id": question_id,
            "user_answer": user_answer,
            "is_correct": is_correct,
            "session_id": session_id,
            "answered_at": datetime.now(KST).isoformat()
        }).execute()

        # 오답 노트 저장
        if not is_correct:
            supabase.table("quiz_incorrect_notes").insert({
                "user_id": user_id,
                "question_id": question_id,
                "reviewed": False,
                "created_at": datetime.now(KST).isoformat()
            }).execute()

        # 피드백 메시지 저장 (여기서는 텍스트만, 해설은 응답 JSON으로 넘김)
        feedback = "✅ 정답입니다!" if is_correct else f"❌ 오답입니다. 정답은 {correct_answer}"
        supabase.table("quiz_messages").insert({
            "session_id": session_id,
            "run_id": run_id,
            "user_id": user_id,
            "role": "ai",
            "kind": "text",
            "payload": json.dumps({"text": feedback}),
        }).execute()

        return JSONResponse({
            "is_correct": is_correct,
            "correct_answer": correct_answer,
            "explanation": explanation,
        })

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})
