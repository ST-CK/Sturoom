from fastapi import FastAPI, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from mangum import Mangum
import os, json, httpx
from io import BytesIO
from datetime import datetime, timedelta, timezone
from typing import Tuple
from openai import OpenAI
from PyPDF2 import PdfReader
from pptx import Presentation
from supabase import create_client, Client

# FastAPI 앱 생성
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "https://sturoom.vercel.app",
        "https://*.vercel.app",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 환경 변수
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")

if not (OPENAI_API_KEY and SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY):
    print("Warning: Some environment variables are missing")
    client = None
    supabase = None
else:
    client = OpenAI(api_key=OPENAI_API_KEY)
    supabase: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)

MAX_TOTAL_CHARS = 18000
KST = timezone(timedelta(hours=9))

# 유틸 함수들
def _safe_cut(text: str, limit: int) -> str:
    return (text or "").strip()[:limit]

async def _download_file(url: str) -> Tuple[str, bytes]:
    async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as http:
        r = await http.get(url)
    if r.status_code != 200:
        raise RuntimeError(f"다운로드 실패({r.status_code})")
    return url.split("/")[-1].lower(), r.content

def _extract_text_from_pdf(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    return "\n".join([page.extract_text() or "" for page in reader.pages])

def _extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    slides = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        slides.append("\n".join(texts))
    return "\n".join(slides)

def _build_prompt(all_text: str, mode: str) -> str:
    if mode == "ox":
        mode_desc = "OX 형식 문제 (정답은 O 또는 X)"
    elif mode == "short":
        mode_desc = "서술형 문제 (짧은 한 문장으로 답변)"
    elif mode == "multiple":
        mode_desc = "객관식 4지선다 문제"
    else:
        mode_desc = "객관식 4지선다, OX, 서술형을 섞은 혼합형 문제"

    return f"""
다음은 강의 자료의 통합 텍스트입니다.
이 내용을 바탕으로 학습 이해도를 평가할 수 있는 {mode_desc} 3문항을 만들어 주세요.

요구사항:
1. JSON 배열로만 출력
2. 각 문항은 question, choices(보기 리스트), answer(정답), explanation(해설)을 포함
3. 보기에는 'A.', 'B.' 같은 접두사를 붙이지 마세요.
4. 출력은 JSON 배열만 반환하세요.

-----
{_safe_cut(all_text, MAX_TOTAL_CHARS)}
-----
"""

# 엔드포인트들
@app.get("/")
def root():
    return {
        "message": "Quiz API",
        "endpoints": ["/session/start", "/from-url", "/attempt"]
    }

@app.post("/session/start")
async def start_quiz_session(req: Request):
    data = await req.json()
    user_id = data.get("user_id")
    room_id = data.get("room_id")
    week_id = data.get("post_id")
    mode = data.get("mode", "mixed")

    if not user_id:
        return JSONResponse(status_code=400, content={"error": "user_id가 필요합니다."})

    try:
        existing = (
            supabase.table("quiz_sessions")
            .select("id")
            .eq("user_id", user_id)
            .eq("lecture_id", room_id)
            .eq("week_id", week_id)
            .limit(1)
            .execute()
        )

        if existing.data and len(existing.data) > 0:
            session_id = existing.data[0]["id"]
        else:
            session = {
                "user_id": user_id,
                "lecture_id": room_id,
                "week_id": week_id,
                "mode": mode,
                "quiz_count": 0,
                "created_at": datetime.now(KST).isoformat(),
            }
            s_res = supabase.table("quiz_sessions").insert(session).execute()
            session_id = s_res.data[0]["id"]

        run = {
            "session_id": session_id,
            "user_id": user_id,
            "lecture_id": room_id,
            "week_id": week_id,
            "mode": mode,
            "started_at": datetime.now(KST).isoformat(),
        }
        r_res = supabase.table("quiz_runs").insert(run).execute()
        run_id = r_res.data[0]["id"]

        return JSONResponse({"session_id": session_id, "run_id": run_id})

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/from-url")
async def generate_quiz_from_url(req: Request):
    data = await req.json()
    file_urls = data.get("file_urls") or []
    mode = (data.get("mode") or "multiple").strip().lower()
    session_id = data.get("session_id")
    run_id = data.get("run_id")
    user_id = data.get("user_id")

    if not file_urls:
        return JSONResponse(status_code=400, content={"error": "파일 URL이 없습니다."})

    aggregated = []
    for file in file_urls:
        try:
            url = file.get("url") if isinstance(file, dict) else file
            fname, blob = await _download_file(url)
            text = _extract_text_from_pdf(blob) if fname.endswith(".pdf") else _extract_text_from_pptx(blob)
            aggregated.append(f"\n### {fname}\n{text}")
        except Exception as e:
            print(f"⚠️ 파일 처리 실패: {url} ({e})")

    all_text = "\n".join(aggregated)
    prompt = _build_prompt(all_text, mode)

    try:
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "너는 교육용 퀴즈를 JSON으로만 반환하는 AI 교사야."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.3,
        )
        quiz_text = resp.choices[0].message.content
        json_str = quiz_text[quiz_text.find("["):quiz_text.rfind("]") + 1]
        quiz_data = json.loads(json_str)
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": f"OpenAI 처리 실패: {str(e)}"})

    try:
        questions = []
        for q in quiz_data:
            answer = q.get("answer") or q.get("correct_answer") or ""
            explanation = q.get("explanation") or ""
            choices = q.get("choices") or q.get("options")
            cleaned_choices = []
            if isinstance(choices, list):
                for c in choices:
                    if isinstance(c, str):
                        c = c.replace("A. ", "").replace("B. ", "").replace("C. ", "").replace("D. ", "").strip()
                    cleaned_choices.append(c)

            questions.append({
                "session_id": session_id,
                "question": q.get("question", "").strip(),
                "choices": cleaned_choices,
                "answer": answer.strip(),
                "explanation": explanation.strip(),
            })

        inserted = supabase.table("quiz_questions").insert(questions).execute()
        supabase.table("quiz_runs").update({"quiz_count": len(inserted.data)}).eq("id", run_id).execute()
        supabase.table("quiz_sessions").update({"quiz_count": len(inserted.data)}).eq("id", session_id).execute()

        supabase.table("quiz_messages").insert({
            "session_id": session_id,
            "run_id": run_id,
            "user_id": user_id,
            "role": "ai",
            "kind": "quiz",
            "payload": json.dumps({"quiz": inserted.data}),
            "order_index": 1
        }).execute()

        return JSONResponse({
            "message": "퀴즈 생성 및 저장 완료",
            "session_id": session_id,
            "run_id": run_id,
            "quiz_count": len(inserted.data),
            "quiz": inserted.data
        })
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/attempt")
async def record_quiz_attempt(req: Request):
    try:
        data = await req.json()
        user_email = data.get("user_email")
        question_id = data.get("question_id")
        user_answer = (data.get("user_answer") or "").strip()
        session_id = data.get("session_id")
        run_id = data.get("run_id")

        if not user_email or not question_id:
            return JSONResponse(status_code=400, content={"error": "user_email 또는 question_id 누락"})

        user_res = supabase.table("profiles").select("id").eq("email", user_email).single().execute()
        if not user_res.data:
            return JSONResponse(status_code=404, content={"error": "해당 이메일 사용자를 찾을 수 없습니다."})
        user_id = user_res.data["id"]

        q_res = supabase.table("quiz_questions").select("answer").eq("id", question_id).single().execute()
        if not q_res.data:
            return JSONResponse(status_code=404, content={"error": "문제를 찾을 수 없습니다."})
        correct_answer = (q_res.data["answer"] or "").strip()
        is_correct = user_answer.lower() == correct_answer.lower()

        supabase.table("quiz_answers").insert({
            "session_id": session_id,
            "question_id": question_id,
            "user_id": user_id,
            "user_answer": user_answer,
            "is_correct": is_correct
        }).execute()

        if not is_correct:
            supabase.table("quiz_incorrect_notes").insert({
                "user_id": user_id,
                "question_id": question_id,
                "reviewed": False
            }).execute()

        order_idx = int(datetime.now().timestamp())
        messages = [
            {
                "session_id": session_id,
                "run_id": run_id,
                "user_id": user_id,
                "role": "user",
                "kind": "text",
                "payload": json.dumps({"text": user_answer}),
                "order_index": order_idx,
            },
            {
                "session_id": session_id,
                "run_id": run_id,
                "user_id": user_id,
                "role": "ai",
                "kind": "text",
                "payload": json.dumps({
                    "text": "✅ 정답입니다!" if is_correct else f"❌ 오답입니다. 정답은 {correct_answer}"
                }),
                "order_index": order_idx + 1,
            },
        ]
        supabase.table("quiz_messages").insert(messages).execute()

        return JSONResponse({
            "message": "정답 판정 완료",
            "is_correct": is_correct,
            "correct_answer": correct_answer,
        }, status_code=200)

    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

handler = Mangum(app, lifespan="off")
